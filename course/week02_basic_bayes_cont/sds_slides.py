"""Shared slide-building module for SDS-branded lecture decks.

Produces slides in the SDS (Chiba Tech, School of Design & Science) style:
- 10 x 5.625 in canvas (matches Google Slides 16:9 default)
- Title + section-break slides use the SDS yellow frame asset from sds_branding.svg
- Agenda slides follow a two-column time-stamped pattern
- Content slides use a dark bg with white text
- Speaker notes embedded directly in each slide's notes pane

Originally built for APS-I at Chiba Tech; reused here for Human and Machine Learning SP26.

Usage pattern (see build_slides_week2.py):

    from sds_slides import SDSDeck

    deck = SDSDeck(course_label="Human and Machine Learning SP26 — Week 2")
    deck.title_slide("Week 2 — Levels of Analysis + Bayes Continued",
                     subtitle="Friday, April 24, 2026", authors="Prof. Austerweil")
    deck.agenda_slide([("Welcome back", "0:00"), ...], highlight="Welcome back")
    deck.section_break("Welcome back", notes="...")
    deck.content_slide("Slide title", [
        ("First line", {"size": 24, "bold": True}),
        ("Second line", {"size": 22, "color": deck.DIM}),
    ], notes="Speaker notes go here — will land directly in the pptx notes pane.")
    ...
    deck.save("week2-slides.pptx")
    deck.export_notes_md("week2-speaker-notes.md", title="Week 2: Speaker Notes")
"""

from __future__ import annotations

import base64
import re
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt


def _estimate_body_height_pt(body, default_size: int, width_in: float) -> float:
    """Estimate rendered height (in points) of a body block, calibrated against
    real PowerPoint rendering of Week 2 HML.

    Empirical factor: sum-of-nominal-sizes × 0.85 matches observed rendered
    height to within ~1 line. Blank-text entries render at ~0.85 of nominal
    font size; text lines render at ~1.0 × size with no extra paragraph
    spacing; wrap advances at ~0.45pt per point of font size.

    Over-estimates by ~5% on purpose so the trim loop is conservative.
    """
    width_pt = width_in * 72.0
    total = 0.0
    for item in body:
        if isinstance(item, str):
            text, opts = item, {}
        else:
            text, opts = item
        size = opts.get("size", default_size)
        if not text:
            total += size * 0.85
            continue
        line_h = size * 1.0
        char_advance_pt = size * 0.45
        chars_per_line = max(1, int(width_pt / char_advance_pt))
        lines_needed = 0
        for seg in text.split("\n") or [text]:
            lines_needed += max(1, -(-len(seg) // chars_per_line))
        total += line_h * lines_needed
    # Calibration factor: real rendered height / raw line-sum. Derived from
    # comparing pre-fix slide 8 (overflowed 4.4in budget by ~1 line → real ~336pt,
    # raw sum ~395pt → factor ~0.85).
    return total * 0.85


# --- Canvas (matches real Day01 deck: 10.00 x 5.625 in) ---
CANVAS_W = Inches(10)
CANVAS_H = Inches(5.625)

# --- Theme colors ---
BG_DARK = RGBColor(0x11, 0x11, 0x11)
BG_SECTION = RGBColor(0x1A, 0x1A, 0x2E)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DIM = RGBColor(0x99, 0x99, 0x99)
TEXT_ACCENT = RGBColor(0x64, 0xB5, 0xF6)
TEXT_RED = RGBColor(0xEF, 0x53, 0x50)
TEXT_GREEN = RGBColor(0x66, 0xBB, 0x6A)
TEXT_ORANGE = RGBColor(0xFF, 0xA7, 0x26)
TEXT_YELLOW = RGBColor(0xFF, 0xEB, 0x3B)

# Section-break text accent matches SDS yellow frame
SDS_YELLOW = RGBColor(0xF5, 0xE4, 0x7A)

_HERE = Path(__file__).parent
_SVG_PATH = _HERE / "sds_branding.svg"
_ASSET_DIR = _HERE / ".sds_assets"


def _extract_sds_assets() -> tuple[Path, Path]:
    """Pull the two PNGs out of Day01_APS-I_Spr2026.svg.

    Returns (frame_path, wordmark_path). The frame is the yellow SDS frame used
    on title + section-break slides; the wordmark is the SDS logo that appears
    in the lower-right of the title slide.
    """
    _ASSET_DIR.mkdir(exist_ok=True)
    frame = _ASSET_DIR / "sds_frame.png"
    wordmark = _ASSET_DIR / "sds_wordmark.png"
    if frame.exists() and wordmark.exists():
        return frame, wordmark

    svg = _SVG_PATH.read_text()
    imgs = re.findall(r'xlink:href="data:image/(?:png|jpeg);base64,([^"]+)"', svg)
    if len(imgs) < 2:
        raise RuntimeError(
            f"Expected 2 inline images in {_SVG_PATH}, found {len(imgs)}."
        )
    frame.write_bytes(base64.b64decode(imgs[0]))
    wordmark.write_bytes(base64.b64decode(imgs[1]))
    return frame, wordmark


class SDSDeck:
    """Wrapper around python-pptx Presentation with SDS styling helpers."""

    BG_DARK = BG_DARK
    BG_SECTION = BG_SECTION
    WHITE = TEXT_WHITE
    DIM = TEXT_DIM
    ACCENT = TEXT_ACCENT
    RED = TEXT_RED
    GREEN = TEXT_GREEN
    ORANGE = TEXT_ORANGE
    YELLOW = TEXT_YELLOW
    SDS_YELLOW = SDS_YELLOW

    def __init__(self, course_label: str = "APS-I Spring 2026"):
        self.prs = Presentation()
        self.prs.slide_width = CANVAS_W
        self.prs.slide_height = CANVAS_H
        self.course_label = course_label
        self._frame_png, self._wordmark_png = _extract_sds_assets()
        self._blank_layout = self.prs.slide_layouts[6]

    # ---------- Internal helpers ----------

    def _set_bg(self, slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_textbox(self, slide, left, top, width, height):
        return slide.shapes.add_textbox(left, top, width, height)

    def _write_lines(self, tf, blocks, default_size=22, default_color=TEXT_WHITE,
                     default_align=PP_ALIGN.LEFT):
        """Write a list of (text, opts) tuples into a text_frame.

        opts keys: size (pt), color (RGBColor), bold (bool), align (PP_ALIGN),
        level (int), space_before (Pt), space_after (Pt), italic (bool).
        """
        tf.word_wrap = True
        first = True
        for item in blocks:
            if isinstance(item, str):
                text, opts = item, {}
            else:
                text, opts = item
            if first:
                p = tf.paragraphs[0]
                p.text = text
                first = False
            else:
                p = tf.add_paragraph()
                p.text = text
            size = opts.get("size", default_size)
            p.font.size = Pt(size)
            p.font.color.rgb = opts.get("color", default_color)
            p.font.bold = opts.get("bold", False)
            p.font.italic = opts.get("italic", False)
            p.alignment = opts.get("align", default_align)
            if "level" in opts:
                p.level = opts["level"]
            if "space_before" in opts:
                p.space_before = opts["space_before"]
            if "space_after" in opts:
                p.space_after = opts["space_after"]

    def _add_notes(self, slide, notes: str | None):
        if not notes:
            return
        # has_notes_slide triggers creation when accessed
        notes_tf = slide.notes_slide.notes_text_frame
        notes_tf.text = notes.strip()

    def _add_frame(self, slide):
        """Overlay the SDS yellow frame across the slide body."""
        slide.shapes.add_picture(
            str(self._frame_png),
            Inches(0.5), Inches(0.5),
            width=Inches(9), height=Inches(4.63),
        )

    def _add_wordmark(self, slide):
        """Place the SDS wordmark in the lower-right (title-slide convention)."""
        slide.shapes.add_picture(
            str(self._wordmark_png),
            Inches(6.67), Inches(3.81),
            width=Inches(2.25), height=Inches(0.75),
        )

    # ---------- Public slide builders ----------

    def title_slide(self, title: str, subtitle: str = "", authors: str = "",
                    notes: str | None = None):
        """Front-of-deck title slide: yellow frame + SDS wordmark + title text."""
        slide = self.prs.slides.add_slide(self._blank_layout)
        self._set_bg(slide, BG_DARK)
        self._add_frame(slide)
        self._add_wordmark(slide)

        tb = self._add_textbox(slide, Inches(1.8), Inches(1.5),
                               Inches(6.4), Inches(2.0))
        blocks = [(title, {"size": 32, "bold": True, "align": PP_ALIGN.LEFT})]
        if authors:
            blocks.append(("", {"size": 10}))
            blocks.append((authors, {"size": 18, "color": TEXT_DIM}))
        if subtitle:
            blocks.append((subtitle, {"size": 16, "color": TEXT_DIM}))
        self._write_lines(tb.text_frame, blocks)

        self._add_notes(slide, notes)
        return slide

    def agenda_slide(self, items: list[tuple[str, str]], highlight: str | None = None,
                     notes: str | None = None):
        """Two-column time-stamped agenda. Repeat before each section.

        items: list of (label, time_str) tuples in chronological order.
        highlight: optional label to render in accent color (shows "where we are").
        """
        slide = self.prs.slides.add_slide(self._blank_layout)
        self._set_bg(slide, BG_DARK)

        # "Agenda" header top-left
        hdr = self._add_textbox(slide, Inches(0.24), Inches(0.18),
                                Inches(2.0), Inches(0.45))
        self._write_lines(hdr.text_frame, [
            ("Agenda", {"size": 22, "bold": True, "color": TEXT_WHITE}),
        ])

        # Two columns: label (centered-ish) + time. Label column is wide
        # enough to fit the longest segment name on one line (e.g.
        # "MP1 Debrief — What Changed?") without wrapping.
        label_left = Inches(3.0)
        label_width = Inches(3.6)
        time_left = Inches(6.7)
        time_width = Inches(0.7)
        row_height = Inches(0.4)
        start_top = Inches(1.1)
        for i, (label, time_str) in enumerate(items):
            top = Inches(1.1 + 0.4 * i)
            is_hl = (highlight is not None and label == highlight)
            color = TEXT_ACCENT if is_hl else TEXT_WHITE
            bold = is_hl
            lb = self._add_textbox(slide, label_left, top, label_width, row_height)
            self._write_lines(lb.text_frame, [
                (label, {"size": 16, "bold": bold, "color": color}),
            ])
            tb = self._add_textbox(slide, time_left, top, time_width, row_height)
            self._write_lines(tb.text_frame, [
                (time_str, {"size": 16, "bold": bold, "color": color}),
            ])

        self._add_notes(slide, notes)
        return slide

    def section_break(self, segment_name: str, lead: str = "",
                      notes: str | None = None):
        """SDS-framed section-break slide: course label, segment name, lead."""
        slide = self.prs.slides.add_slide(self._blank_layout)
        self._set_bg(slide, BG_DARK)
        self._add_frame(slide)

        # Course label at top (inside the yellow frame)
        top_tb = self._add_textbox(slide, Inches(1.8), Inches(0.55),
                                   Inches(6.4), Inches(0.4))
        self._write_lines(top_tb.text_frame, [
            (self.course_label, {"size": 14, "color": TEXT_DIM,
                                  "align": PP_ALIGN.CENTER}),
        ])

        # Segment name centered in frame (horizontally + vertically)
        mid_tb = self._add_textbox(slide, Inches(1.8), Inches(1.9),
                                   Inches(6.4), Inches(1.8))
        self._write_lines(mid_tb.text_frame, [
            (segment_name, {"size": 32, "bold": True, "color": TEXT_WHITE,
                            "align": PP_ALIGN.CENTER}),
        ])

        # Lead name centered below the segment name
        if lead:
            lead_tb = self._add_textbox(slide, Inches(1.8), Inches(3.6),
                                        Inches(6.4), Inches(0.6))
            self._write_lines(lead_tb.text_frame, [
                (lead, {"size": 20, "color": TEXT_DIM,
                        "align": PP_ALIGN.CENTER}),
            ])

        self._add_notes(slide, notes)
        return slide

    def content_slide(self, title: str, body, notes: str | None = None,
                      title_color=TEXT_WHITE):
        """Standard content slide: title bar + body block.

        body: list of (text, opts) tuples — see _write_lines for opts.
        """
        slide = self.prs.slides.add_slide(self._blank_layout)
        self._set_bg(slide, BG_DARK)

        # Title bar at top
        title_tb = self._add_textbox(slide, Inches(0.6), Inches(0.3),
                                     Inches(8.8), Inches(0.56))
        self._write_lines(title_tb.text_frame, [
            (title, {"size": 22, "bold": True, "color": title_color}),
        ])

        # Body region. MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE is unreliable — LibreOffice
        # pre-shrinks in the PDF preview but PowerPoint does not shrink until the
        # file is edited. Instead, trim lines from the TOP at build time so the
        # punchline (usually last) is preserved. Preserves font size (readable on
        # projector) at the cost of dropping earlier build-up context.
        body_in = (0.6, 1.0, 8.8, 4.4)  # left, top, width, height
        max_body_pt = body_in[3] * 72.0
        trimmed_body = list(body)
        while (len(trimmed_body) > 1
               and _estimate_body_height_pt(trimmed_body, 18, body_in[2]) > max_body_pt):
            trimmed_body.pop(0)
        dropped = len(body) - len(trimmed_body)
        if dropped:
            print(f"[content_slide] '{title}': trimmed {dropped} line(s) from top to fit")
        body_tb = self._add_textbox(slide, Inches(body_in[0]), Inches(body_in[1]),
                                    Inches(body_in[2]), Inches(body_in[3]))
        body_tb.text_frame.word_wrap = True
        self._write_lines(body_tb.text_frame, trimmed_body, default_size=18)

        self._add_notes(slide, notes)
        return slide

    def content_image_slide(self, title: str, image_path: str,
                            caption=None, notes: str | None = None,
                            title_color=TEXT_WHITE,
                            image_width_in: float = 7.5,
                            image_top_in: float = 1.0):
        """Slide with a title bar and a centered image below.

        image_path: absolute or relative path to a PNG/JPG.
        caption: optional list of (text, opts) tuples rendered below the image.
        image_width_in: image width in inches (height scales from the source file).
        image_top_in: vertical position of the image in inches from slide top.
        """
        slide = self.prs.slides.add_slide(self._blank_layout)
        self._set_bg(slide, BG_DARK)

        title_tb = self._add_textbox(slide, Inches(0.6), Inches(0.3),
                                     Inches(8.8), Inches(0.56))
        self._write_lines(title_tb.text_frame, [
            (title, {"size": 22, "bold": True, "color": title_color}),
        ])

        # Center the image horizontally
        left = Inches((10 - image_width_in) / 2)
        slide.shapes.add_picture(str(image_path), left, Inches(image_top_in),
                                 width=Inches(image_width_in))

        if caption:
            cap_w_in = 8.8
            cap_h_in = 0.65
            cap_top_in = 5.625 - cap_h_in - 0.02  # pin to bottom with small margin
            max_cap_pt = cap_h_in * 72.0
            trimmed_cap = list(caption)
            while (len(trimmed_cap) > 1
                   and _estimate_body_height_pt(trimmed_cap, 16, cap_w_in) > max_cap_pt):
                trimmed_cap.pop(0)
            dropped = len(caption) - len(trimmed_cap)
            if dropped:
                print(f"[content_image_slide] '{title}': trimmed {dropped} caption line(s)")
            cap_tb = self._add_textbox(slide, Inches(0.6), Inches(cap_top_in),
                                       Inches(cap_w_in), Inches(cap_h_in))
            cap_tb.text_frame.word_wrap = True
            self._write_lines(cap_tb.text_frame, trimmed_cap, default_size=16)

        self._add_notes(slide, notes)
        return slide

    def build_slides(self, title: str, steps: list, notes: str | list | None = None,
                     title_color=TEXT_WHITE):
        """Emit a sequence of content slides that share a title and build up a body
        block line-by-line.

        Each entry in `steps` is itself a `body` list (same shape as `content_slide`'s
        body). Slide N renders the concatenation of `steps[0] + steps[1] + ... + steps[N]`,
        so advancing reveals the next batch of lines on top of what's already there.

        `notes` may be a single string (applied to the last slide only; earlier slides
        get no notes), or a list of strings the same length as `steps` (one per slide).

        Use for math build-ups: show the formula, then name each piece, then plug
        in numbers, etc.
        """
        if isinstance(notes, str) or notes is None:
            notes_list = [None] * (len(steps) - 1) + [notes]
        else:
            assert len(notes) == len(steps), (
                f"notes list length {len(notes)} must match steps length {len(steps)}"
            )
            notes_list = notes
        accumulated = []
        slides = []
        for step, step_notes in zip(steps, notes_list):
            accumulated = accumulated + step
            slide = self.content_slide(
                title=title, body=list(accumulated),
                notes=step_notes, title_color=title_color,
            )
            slides.append(slide)
        return slides

    def break_slide(self, text: str = "Break", notes: str | None = None):
        """Framed 'break' slide for mid-lecture pauses."""
        slide = self.prs.slides.add_slide(self._blank_layout)
        self._set_bg(slide, BG_DARK)
        self._add_frame(slide)

        top_tb = self._add_textbox(slide, Inches(1.8), Inches(0.55),
                                   Inches(6.4), Inches(0.4))
        self._write_lines(top_tb.text_frame, [
            (self.course_label, {"size": 14, "color": TEXT_DIM,
                                  "align": PP_ALIGN.CENTER}),
        ])

        mid_tb = self._add_textbox(slide, Inches(1.8), Inches(2.0),
                                   Inches(6.4), Inches(1.5))
        self._write_lines(mid_tb.text_frame, [
            (text, {"size": 36, "bold": True, "color": TEXT_WHITE,
                    "align": PP_ALIGN.CENTER}),
        ])

        self._add_notes(slide, notes)
        return slide

    # ---------- Output ----------

    def save(self, path: str):
        self.prs.save(path)
        print(f"Saved to {path}")
        print(f"Total slides: {len(self.prs.slides)}")

    def render_previews(self, pptx_path: str, out_dir: str, dpi: int = 150):
        """Render each slide to a PNG for visual-feedback review.

        pptx → pdf via LibreOffice headless → per-page PNG via pdftoppm (poppler).
        Writes slide01.png, slide02.png, … into out_dir. Requires `soffice` on PATH
        and `pdf2image` + `poppler-utils` installed.
        """
        import shutil
        import subprocess
        import tempfile

        from pdf2image import convert_from_path

        pptx = Path(pptx_path).resolve()
        out = Path(out_dir).resolve()
        out.mkdir(parents=True, exist_ok=True)
        for stale in out.glob("slide*.png"):
            stale.unlink()

        soffice = shutil.which("soffice") or shutil.which("libreoffice")
        if not soffice:
            raise RuntimeError("LibreOffice (soffice) not found on PATH")

        with tempfile.TemporaryDirectory() as td:
            subprocess.run(
                [soffice, "--headless", "--convert-to", "pdf",
                 "--outdir", td, str(pptx)],
                check=True, capture_output=True,
            )
            pdf = Path(td) / f"{pptx.stem}.pdf"
            if not pdf.exists():
                raise RuntimeError(f"LibreOffice did not produce {pdf}")
            images = convert_from_path(str(pdf), dpi=dpi)

        for i, img in enumerate(images, 1):
            img.save(out / f"slide{i:02d}.png", "PNG")
        print(f"Rendered {len(images)} preview PNGs to {out}")

    def export_notes_md(self, path: str, title: str = "Speaker Notes"):
        """Dump every slide's speaker notes to a readable markdown file.

        This is the companion artifact for co-lecturers. Slides without notes
        are listed with a placeholder so gaps are visible.
        """
        lines = [f"# {title}", ""]
        for i, slide in enumerate(self.prs.slides, 1):
            slide_title = _slide_heading(slide) or f"Slide {i}"
            lines.append(f"## Slide {i}: {slide_title}")
            lines.append("")
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
            else:
                notes = ""
            lines.append(notes if notes else "_(no speaker notes)_")
            lines.append("")
        Path(path).write_text("\n".join(lines))
        print(f"Wrote notes markdown to {path}")


def _slide_heading(slide) -> str | None:
    """Pick the best text to use as the slide's heading in the notes markdown.

    Prefers the largest-font non-empty text on the slide (typically the title or
    the centered segment name on a section break). Falls back to the first
    non-empty text. Skips the course-label banner so section-break slides show
    their segment name instead of the repeated header text.
    """
    candidates = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        first_line = text.split("\n")[0]
        max_size = 0
        for para in shape.text_frame.paragraphs:
            if not para.text.strip():
                continue
            if para.font.size:
                max_size = max(max_size, para.font.size.pt)
            for run in para.runs:
                if run.font.size:
                    max_size = max(max_size, run.font.size.pt)
            break  # only first paragraph
        # Skip the course-label banner (14pt) so the segment name wins on section breaks.
        if max_size and max_size < 18:
            continue
        # Shape vertical position — title boxes sit near the top of the slide;
        # body content sits lower. When two shapes have similar font sizes,
        # prefer the upper one (the title).
        top_emu = shape.top if shape.top is not None else 10_000_000
        candidates.append((top_emu, -max_size, first_line[:120]))
    if not candidates:
        return None
    # Sort: topmost shape first; among shapes at the same top, largest font wins.
    candidates.sort(key=lambda x: (x[0], x[1]))
    return candidates[0][2]
