"""Build week01.pptx for HML SP26 in the APS-I template style.

Visual language:
  - 10 x 5.625 in (16:9 widescreen)
  - Black background everywhere
  - Title slide: yellow ticket frame around centered title; SDS logo bottom-right
  - Outline (agenda) slide: small "Agenda" eyebrow top-left, SDS mark top-right,
    two-column block of {item title | time}
  - Section divider: mint ticket frame, "Spring 2026 Day 1" centered above frame
    (in dark text on the mint), big block-title centered inside
  - Content slide: small white title top-left, content below; SDS mark top-right
  - Fonts: Helvetica Neue for chrome/headers, Calibri for content

Run: python build_pptx.py  (writes week01.pptx in the same dir)
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = Path(__file__).parent
IMG = HERE / "images"

# ---------- color palette (from APS slides 4-5) ----------
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x99, 0x99, 0x99)
DIM = RGBColor(0xCC, 0xCC, 0xCC)
BLUE = RGBColor(0x64, 0xB5, 0xF6)
RED = RGBColor(0xEF, 0x53, 0x50)
GREEN = RGBColor(0x66, 0xBB, 0x6A)
YELLOW = RGBColor(0xFF, 0xE8, 0x6D)
MINT = RGBColor(0x7F, 0xE5, 0xBC)

DAY_LABEL = "Human and Machine Learning  ·  Week 1"


def set_bg(slide, color=BLACK):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625)
    )
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg


def add_text(
    slide, x, y, w, h, text, *,
    size=14, bold=False, color=WHITE, font="Calibri", align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if isinstance(text, str):
        runs = [(text, {})]
    else:
        runs = text  # list of (text, {style overrides})
    first = True
    for txt, ov in runs:
        if first:
            r = p.add_run()
            first = False
        else:
            # newline -> new paragraph (Calibri behavior)
            p = tf.add_paragraph()
            p.alignment = align
            r = p.add_run()
        r.text = txt
        f = r.font
        f.name = ov.get("font", font)
        f.size = Pt(ov.get("size", size))
        f.bold = ov.get("bold", bold)
        col = ov.get("color", color)
        f.color.rgb = col
    return tb


def add_corner_marks(slide, day_text=DAY_LABEL):
    # tiny header line top-left + day label centered top + page mark bottom-right will be page numbers
    add_text(slide, 0.30, 0.20, 4.0, 0.30, day_text,
             size=9, bold=True, color=MUTED, font="Helvetica Neue")


def add_footer_pageno(slide, n):
    add_text(slide, 9.55, 5.30, 0.40, 0.25, str(n),
             size=9, color=MUTED, font="Helvetica Neue", align=PP_ALIGN.RIGHT)


# ---------- slide builders ----------
def slide_title(prs, title, subtitle, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide)
    # yellow ticket frame
    pic = slide.shapes.add_picture(str(IMG / "frame_yellow.png"),
                                    Inches(0.50), Inches(0.50),
                                    Inches(9.00), Inches(4.63))
    # title
    add_text(slide, 1.0, 1.55, 8.0, 1.7, title,
             size=40, bold=True, color=WHITE, font="Helvetica Neue",
             align=PP_ALIGN.CENTER)
    # subtitle
    add_text(slide, 1.0, 3.2, 8.0, 0.4, subtitle,
             size=18, color=MUTED, font="Helvetica Neue",
             align=PP_ALIGN.CENTER)
    add_text(slide, 1.0, 3.7, 8.0, 0.4, meta,
             size=14, color=MUTED, font="Helvetica Neue",
             align=PP_ALIGN.CENTER)
    # SDS logo bottom-right inside the frame
    slide.shapes.add_picture(str(IMG / "sds_logo_white.png"),
                              Inches(7.0), Inches(4.30),
                              Inches(2.0), Inches(0.65))
    return slide


def slide_agenda(prs, items, page=2):
    """items: list of (title, time_str)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    # eyebrow
    add_text(slide, 0.30, 0.18, 2.0, 0.45, "Outline",
             size=15, bold=True, color=WHITE, font="Helvetica Neue")
    # day text top-right (right-aligned)
    add_text(slide, 5.6, 0.20, 4.1, 0.30, DAY_LABEL,
             size=9, bold=True, color=MUTED, font="Helvetica Neue",
             align=PP_ALIGN.RIGHT)
    # two-column layout for items, vertically centered-ish
    n = len(items)
    row_h = 0.42
    block_h = n * row_h
    top = (5.625 - block_h) / 2 - 0.1
    title_x = 2.80
    title_w = 3.6
    time_x = 6.55
    time_w = 2.5
    for i, (title, time_str) in enumerate(items):
        y = top + i * row_h
        add_text(slide, title_x, y, title_w, row_h, title,
                 size=13, bold=True, color=WHITE, font="Helvetica Neue",
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, time_x, y, time_w, row_h, time_str,
                 size=10, color=MUTED, font="Helvetica Neue",
                 anchor=MSO_ANCHOR.MIDDLE)
    add_footer_pageno(slide, page)
    return slide


def slide_section(prs, title, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    # mint ticket frame
    slide.shapes.add_picture(str(IMG / "frame_mint.png"),
                              Inches(0.50), Inches(0.50),
                              Inches(9.00), Inches(4.63))
    # day label centered above the frame's inner dark area (on the mint top bar)
    add_text(slide, 1.80, 0.55, 6.39, 0.37, DAY_LABEL,
             size=9, bold=True, color=BLACK, font="Helvetica Neue",
             align=PP_ALIGN.CENTER)
    # big section title centered inside the frame
    add_text(slide, 1.0, 2.05, 8.0, 1.5, title,
             size=42, bold=True, color=WHITE, font="Helvetica Neue",
             align=PP_ALIGN.CENTER)
    add_footer_pageno(slide, page)
    return slide


def slide_content(prs, header, body, page):
    """body: list of (text, style-dict). Empty text is a blank line.
    style keys: size, bold, color, font, indent (int 0/1/2)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    # small white header
    add_text(slide, 0.60, 0.30, 8.78, 0.55, header,
             size=22, bold=True, color=WHITE, font="Helvetica Neue")
    # day mark right
    add_text(slide, 5.6, 0.20, 4.1, 0.25, DAY_LABEL,
             size=9, bold=True, color=MUTED, font="Helvetica Neue",
             align=PP_ALIGN.RIGHT)
    # content
    tb = slide.shapes.add_textbox(Inches(0.60), Inches(1.05),
                                   Inches(8.78), Inches(4.0))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for txt, st in body:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        indent = st.get("indent", 0)
        p.level = indent
        if not txt:
            # blank line spacer — keep small
            p.space_after = Pt(0)
            continue
        r = p.add_run()
        r.text = txt
        f = r.font
        f.name = st.get("font", "Calibri")
        f.size = Pt(st.get("size", 16))
        f.bold = st.get("bold", False)
        f.italic = st.get("italic", False)
        f.color.rgb = st.get("color", DIM)
        p.space_after = Pt(st.get("space_after", 6))
        if st.get("align"):
            p.alignment = st["align"]
    add_footer_pageno(slide, page)
    return slide


def slide_centered(prs, big, small, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, 0.60, 2.20, 8.78, 1.0, big,
             size=34, bold=True, color=BLUE, font="Helvetica Neue",
             align=PP_ALIGN.CENTER)
    if small:
        add_text(slide, 0.60, 3.30, 8.78, 1.0, small,
                 size=16, color=MUTED, font="Helvetica Neue",
                 align=PP_ALIGN.CENTER)
    add_footer_pageno(slide, page)
    return slide


def slide_grading(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, 0.60, 0.30, 8.78, 0.55, "Admin  —  grading",
             size=22, bold=True, color=WHITE, font="Helvetica Neue")
    add_text(slide, 5.6, 0.20, 4.1, 0.25, DAY_LABEL,
             size=9, bold=True, color=MUTED, font="Helvetica Neue",
             align=PP_ALIGN.RIGHT)
    rows = [
        ("Final project", "50%",
         "proposal 5%  ·  talk 7.5%  ·  paper 37.5%"),
        ("Programming assignments (4)", "30%",
         "Clusters 7.5%  ·  Gen 7.5%  ·  MC 10.5%  ·  RL 4.5%"),
        ("Weekly written reflections", "15%",
         "~200 words  ·  8 of 13  ·  pass / fail"),
        ("Participation", "5%", ""),
        ("Quizzes", "0%", "self-check only"),
    ]
    top = 1.10
    row_h = 0.72
    for i, (name, pct, note) in enumerate(rows):
        y = top + i * row_h
        add_text(slide, 0.80, y, 5.4, 0.40, name,
                 size=16, bold=True, color=WHITE, font="Helvetica Neue",
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, 6.20, y, 1.2, 0.40, pct,
                 size=18, bold=True, color=BLUE, font="Helvetica Neue",
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
        if note:
            add_text(slide, 0.80, y + 0.40, 8.4, 0.28, note,
                     size=12, color=MUTED, font="Helvetica Neue",
                     anchor=MSO_ANCHOR.TOP)
    add_text(slide, 0.60, 5.00, 8.78, 0.30,
             "All four assignments are completed in GenJAX.",
             size=13, color=DIM, font="Helvetica Neue")
    add_footer_pageno(slide, page)
    return slide


def slide_image_right(prs, header, body, image_path, page,
                      img_w=4.4, img_h=3.6, img_x=5.2, img_y=1.05):
    """header + bullet body left, image right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, 0.60, 0.30, 8.78, 0.55, header,
             size=22, bold=True, color=WHITE, font="Helvetica Neue")
    add_text(slide, 5.6, 0.20, 4.1, 0.25, DAY_LABEL,
             size=9, bold=True, color=MUTED, font="Helvetica Neue",
             align=PP_ALIGN.RIGHT)
    # body left
    tb = slide.shapes.add_textbox(Inches(0.60), Inches(1.05),
                                   Inches(4.5), Inches(4.0))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for txt, st in body:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        if not txt:
            continue
        r = p.add_run()
        r.text = txt
        f = r.font
        f.name = st.get("font", "Calibri")
        f.size = Pt(st.get("size", 15))
        f.bold = st.get("bold", False)
        f.color.rgb = st.get("color", DIM)
        p.space_after = Pt(st.get("space_after", 6))
    slide.shapes.add_picture(str(image_path),
                              Inches(img_x), Inches(img_y),
                              Inches(img_w), Inches(img_h))
    add_footer_pageno(slide, page)
    return slide


# ---------- content ----------
AGENDA = [
    ("Welcome",                 "Block 1  ·  5 min"),
    ("Why this course exists",  "Block 2  ·  25 min"),
    ("Basic Bayes",             "Block 3  ·  25 min"),
    ("Bayes in action",         "Block 4  ·  35 min"),
    ("GenJAX, admin, homework", "Block 5  ·  20 min"),
]


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # 1. title
    slide_title(prs,
                "Human and Machine Learning",
                "Week 1  —  Introduction and Basic Bayes",
                "Spring 2026  ·  Apr 17  ·  Prof. Joseph Austerweil")

    # 2. agenda
    slide_agenda(prs, AGENDA, page=2)

    # 3. block 1 divider
    slide_section(prs, "Welcome", page=3)

    # 4. today
    slide_content(prs, "Today",
                  [
                      ("•  Who you are, what you're hoping to get out of the course",
                       {"size": 16, "color": DIM}),
                      ("•  Why this course exists  —  inductive inference under uncertainty",
                       {"size": 16, "color": DIM}),
                      ("•  Basic probability and Bayes' rule",
                       {"size": 16, "color": DIM}),
                      ("•  Two worked examples: sick friend, Kahneman–Tversky cab problem",
                       {"size": 16, "color": DIM}),
                      ("•  What is GenJAX, and how will we use it?",
                       {"size": 16, "color": DIM}),
                      ("•  Logistics and homework",
                       {"size": 16, "color": DIM}),
                  ], page=4)

    # 5. block 2 divider
    slide_section(prs, "Why This Course Exists", page=5)

    # 6. deduction vs induction
    slide_content(prs, "Deduction vs. induction",
                  [
                      ("Deduction", {"size": 18, "bold": True, "color": WHITE}),
                      ("1 + 4 = ?", {"size": 22, "color": BLUE,
                                     "font": "Cambria Math"}),
                      ("Preserves truth. One answer.",
                       {"size": 14, "color": MUTED}),
                      ("", {}),
                      ("Induction", {"size": 18, "bold": True, "color": WHITE}),
                      ("? + ? = 5", {"size": 22, "color": RED,
                                     "font": "Cambria Math"}),
                      ("Underdetermined. Many plausible answers.",
                       {"size": 14, "color": MUTED}),
                      ("", {}),
                      ("The problems people excel at — where we outperform machines — are inductive.",
                       {"size": 16, "color": WHITE, "bold": True}),
                      ("They feel easy. They are not.",
                       {"size": 14, "color": MUTED, "italic": True}),
                  ], page=6)

    # 7. centered: three inductive problems
    slide_centered(prs, "Three inductive problems",
                   "(... that the mind solves constantly)", page=7)

    # 8. perception
    slide_image_right(prs, "1.  Perception",
                      [
                          ("Which square is darker — A or B?",
                           {"size": 18, "bold": True, "color": WHITE,
                            "space_after": 12}),
                          ("The visual system solves",
                           {"size": 14, "color": DIM}),
                          ("square color + shadow = intensity",
                           {"size": 14, "color": BLUE, "font": "Consolas"}),
                          ("for square color, using priors about how retinal images are generated.",
                           {"size": 14, "color": DIM, "space_after": 12}),
                          ("You are literally looking at an inductive inference right now.",
                           {"size": 14, "bold": True, "color": WHITE,
                            "italic": True}),
                      ],
                      IMG / "checkershadow.jpg", page=8,
                      img_x=5.3, img_y=1.05, img_w=4.4, img_h=3.5)

    # 9. mental states (setup)
    slide_content(prs, "2.  Mental states",
                  [
                      ("Data:  other people's behavior (motion, speech, gaze).",
                       {"size": 16, "color": DIM}),
                      ("Hypotheses:  their goals and beliefs.",
                       {"size": 16, "color": DIM, "space_after": 14}),
                      ("Heider & Simmel (1944) — animated geometric shapes.",
                       {"size": 14, "color": MUTED, "italic": True,
                        "space_after": 18}),
                      ("Watch ~90 seconds, then describe what happened:",
                       {"size": 16, "bold": True, "color": WHITE}),
                      ("https://www.youtube.com/watch?v=VTNmLt7QX8E",
                       {"size": 14, "color": BLUE}),
                  ], page=9)

    # 10. mental states (debrief)
    slide_content(prs, "2.  Mental states  —  debrief",
                  [
                      ("Everyone spontaneously narrates the shapes as agents with intentions, grudges, fears.",
                       {"size": 15, "color": DIM, "space_after": 8}),
                      ("Those intentions were never in the video.  Shapes moved on a plane.",
                       {"size": 15, "bold": True, "color": WHITE,
                        "space_after": 18}),
                      ("•  Data:  2D motion trajectories.",
                       {"size": 15, "color": DIM}),
                      ("•  Hidden:  goals, beliefs, relationships.",
                       {"size": 15, "color": DIM, "space_after": 14}),
                      ("Your mind solved the inverse problem — behavior → mental states — without you noticing.",
                       {"size": 14, "color": MUTED, "italic": True}),
                  ], page=10)

    # 11. word learning
    slide_content(prs, "3.  Word learning",
                  [
                      ("You're in the Australian outback. Someone points at a hopping animal and says \"jumbuck.\"",
                       {"size": 15, "color": DIM, "space_after": 8}),
                      ("What does jumbuck mean?",
                       {"size": 16, "bold": True, "color": WHITE,
                        "space_after": 8}),
                      ("•  the animal itself?", {"size": 14, "color": DIM}),
                      ("•  undetached kangaroo-parts?", {"size": 14, "color": DIM}),
                      ("•  kangaroo temporal-stages?", {"size": 14, "color": DIM}),
                      ("•  any mammal?", {"size": 14, "color": DIM}),
                      ("•  dinner?", {"size": 14, "color": DIM, "space_after": 12}),
                      ("All consistent with the data. Children pick one — fast.",
                       {"size": 14, "color": DIM}),
                      ("They bring strong prior expectations about what words mean.",
                       {"size": 14, "color": MUTED, "italic": True}),
                  ], page=11)

    # 12. check-in
    slide_content(prs, "Check-in",
                  [
                      ("What's the common structure across these three problems?",
                       {"size": 18, "bold": True, "color": WHITE,
                        "space_after": 16}),
                      ("•  Hidden variable we care about:  h",
                       {"size": 15, "color": DIM}),
                      ("•  Noisy / sparse data we observe:  d",
                       {"size": 15, "color": DIM}),
                      ("•  Prior knowledge we bring to bear",
                       {"size": 15, "color": DIM, "space_after": 16}),
                      ("P(h | d)  ∝  P(d | h)  ·  P(h)",
                       {"size": 22, "color": BLUE, "bold": True,
                        "font": "Cambria Math"}),
                      ("(We'll unpack every symbol in the next hour.)",
                       {"size": 13, "color": MUTED, "italic": True}),
                  ], page=12)

    # 13. block 3 divider
    slide_section(prs, "Basic Bayes", page=13)

    # 14. chibany is hungry
    slide_content(prs, "Chibany is hungry",
                  [
                      ("Chibany, mascot of Chiba Tech, gets two bento offerings per day:",
                       {"size": 16, "color": DIM, "space_after": 6}),
                      ("one for lunch, one for dinner.  Each is either Hamburger (H) or Tonkatsu (T).",
                       {"size": 16, "color": DIM, "space_after": 14}),
                      ("Outcome space:   \u03A9 = {HH, HT, TH, TT}     |\u03A9| = 4",
                       {"size": 17, "color": WHITE, "bold": True,
                        "font": "Cambria Math", "space_after": 10}),
                      ("•  Event:  any subset of \u03A9",
                       {"size": 15, "color": DIM}),
                      ("       e.g., \"at least one tonkatsu\" = {HT, TH, TT}",
                       {"size": 14, "color": MUTED}),
                      ("•  Event space:  the set of all events (powerset of \u03A9).   2\u2074 = 16 events.",
                       {"size": 15, "color": DIM, "space_after": 12}),
                      ("(This example is from the course textbook, Ch 2.)",
                       {"size": 12, "color": MUTED, "italic": True}),
                  ], page=14)

    # 15. probability is counting
    slide_content(prs, "Probability is counting",
                  [
                      ("Chibany wants tonkatsu.  What's the probability of at least one T?",
                       {"size": 16, "color": DIM, "space_after": 12}),
                      ("P(A)  =  |A|  /  |\u03A9|",
                       {"size": 22, "color": WHITE, "bold": True,
                        "font": "Cambria Math", "space_after": 8}),
                      ("A = {HT, TH, TT},   |A| = 3,   |\u03A9| = 4",
                       {"size": 16, "color": DIM, "font": "Cambria Math",
                        "space_after": 8}),
                      ("P(at least one T)  =  3/4  =  0.75",
                       {"size": 20, "color": BLUE, "bold": True,
                        "font": "Cambria Math", "space_after": 16}),
                      ("Count the outcomes in the event.  Divide by total outcomes.  That's it.",
                       {"size": 15, "color": DIM}),
                  ], page=15)

    # 16. random variables
    slide_content(prs, "Random variables",
                  [
                      ("Chibany wants to know:  how much tonkatsu?",
                       {"size": 16, "color": DIM, "space_after": 10}),
                      ("A random variable is a function from outcomes to numbers.",
                       {"size": 16, "color": DIM, "space_after": 10}),
                      ("f : \u03A9 \u2192 {0, 1, 2}     \"count the tonkatsus\"",
                       {"size": 17, "color": WHITE, "bold": True,
                        "font": "Cambria Math", "space_after": 8}),
                      ("f(HH) = 0,    f(HT) = 1,    f(TH) = 1,    f(TT) = 2",
                       {"size": 16, "color": DIM, "font": "Cambria Math",
                        "space_after": 12}),
                      ("P(f = 1)  =  |{HT, TH}| / |\u03A9|  =  2/4  =  1/2",
                       {"size": 18, "color": BLUE, "bold": True,
                        "font": "Cambria Math", "space_after": 12}),
                      ("The word random is about the input, not the mapping.",
                       {"size": 13, "color": MUTED, "italic": True}),
                  ], page=16)

    # 17. joint probability
    slide_content(prs, "Joint probability",
                  [
                      ("What's the chance of hamburger for lunch and tonkatsu for dinner?",
                       {"size": 16, "color": DIM, "space_after": 12}),
                      ("P(lunch = H,  dinner = T)  =  P({HT})  =  1/4",
                       {"size": 20, "color": WHITE, "bold": True,
                        "font": "Cambria Math", "space_after": 16}),
                      ("Joint = \"both of these at once.\"",
                       {"size": 15, "color": DIM}),
                  ], page=17)

    # 18. conditional probability — Tanaka-san
    slide_content(prs, "Conditional probability  \u2014  Tanaka-san visits",
                  [
                      ("Tanaka-san tells Chibany: \"Tomorrow you'll get at least one tonkatsu!\"",
                       {"size": 16, "color": DIM, "space_after": 6}),
                      ("Chibany asks:  what's the probability dinner is also tonkatsu?",
                       {"size": 16, "color": DIM, "space_after": 6}),
                      ("Tanaka-san says 1/2.  Chibany disagrees.",
                       {"size": 16, "bold": True, "color": WHITE,
                        "space_after": 14}),
                      ("Conditioning = restricting the outcome space.",
                       {"size": 15, "color": DIM, "space_after": 6}),
                      ("Cross out HH.  New universe:  {HT, TH, TT}  (3 outcomes).",
                       {"size": 15, "color": DIM}),
                      ("Of those, 2 have dinner = T.",
                       {"size": 15, "color": DIM, "space_after": 10}),
                      ("P(dinner = T  |  \u22651 T)  =  2/3",
                       {"size": 22, "color": BLUE, "bold": True,
                        "font": "Cambria Math"}),
                  ], page=18)

    # 19. dependence vs independence
    slide_content(prs, "Dependence vs. independence",
                  [
                      ("Tanaka-san protests:  \"But what if I just tell you the first meal is T?\"",
                       {"size": 16, "color": DIM, "space_after": 8}),
                      ("P(dinner = T  |  lunch = T)  =  |{TT}| / |{TH, TT}|  =  1/2",
                       {"size": 17, "color": WHITE, "bold": True,
                        "font": "Cambria Math", "space_after": 14}),
                      ("That's the same as the unconditional P(dinner = T) = 1/2.  No change!",
                       {"size": 15, "color": DIM, "space_after": 12}),
                      ("Dependent:   P(A | B) \u2260 P(A)    —  learning B changes your belief about A.",
                       {"size": 15, "color": DIM}),
                      ("Independent:  P(A | B) = P(A)    —  learning B tells you nothing new.",
                       {"size": 15, "color": DIM, "space_after": 12}),
                      ("\"At least one T\" and \"lunch = T\" give different information about dinner!",
                       {"size": 14, "color": MUTED, "italic": True}),
                  ], page=19)

    # 20. marginalization
    slide_content(prs, "Marginalization",
                  [
                      ("Chibany's dinner student is sick.  Only one meal today:  \u03A9\u2081 = {H, T}.",
                       {"size": 15, "color": DIM, "space_after": 6}),
                      ("P(lunch = T) = 1/2.",
                       {"size": 15, "color": DIM, "space_after": 10}),
                      ("Next day, both meals are back.  What's P(lunch = T) now?",
                       {"size": 15, "color": DIM, "space_after": 6}),
                      ("P(lunch = T)  =  P(TH) + P(TT)  =  1/4 + 1/4  =  1/2     Same!",
                       {"size": 16, "color": WHITE, "font": "Cambria Math",
                        "space_after": 14}),
                      ("P(X)  =  \u03A3_c  P(X, C = c)",
                       {"size": 22, "color": BLUE, "bold": True,
                        "font": "Cambria Math", "space_after": 10}),
                      ("Sum the joint over the values of whatever you want to get rid of.",
                       {"size": 14, "color": MUTED, "italic": True}),
                  ], page=20)

    # 21. the ratio definition + Bayes' rule
    slide_content(prs, "The ratio definition  \u2192  Bayes' rule",
                  [
                      ("P(A | B)  =  P(A \u2229 B)  /  P(B)",
                       {"size": 20, "color": WHITE, "bold": True,
                        "font": "Cambria Math", "space_after": 8}),
                      ("Chibany checks:  P(dinner=T | \u22651T) = (2/4) / (3/4) = 2/3   \u2713",
                       {"size": 14, "color": DIM, "font": "Cambria Math",
                        "space_after": 16}),
                      ("From the product rule, both directions:",
                       {"size": 15, "color": DIM, "space_after": 6}),
                      ("P(A | B) \u00b7 P(B)  =  P(A, B)  =  P(B | A) \u00b7 P(A)",
                       {"size": 17, "color": WHITE, "font": "Cambria Math",
                        "space_after": 8}),
                      ("Divide:",
                       {"size": 15, "color": DIM, "space_after": 6}),
                      ("P(A | B)  =  P(B | A) \u00b7 P(A)  /  P(B)",
                       {"size": 24, "bold": True, "color": BLUE,
                        "font": "Cambria Math", "space_after": 14}),
                      ("That's it. The rest of the course is understanding what this means.",
                       {"size": 14, "color": MUTED, "italic": True}),
                  ], page=21)

    # 22. check-in
    slide_content(prs, "Check-in",
                  [
                      ("For the problem I'm about to give you — the cab problem —",
                       {"size": 17, "bold": True, "color": WHITE,
                        "space_after": 14}),
                      ("•  What's d (the data, the observation)?",
                       {"size": 16, "color": DIM}),
                      ("•  What's h (the hidden variable)?",
                       {"size": 16, "color": DIM}),
                  ], page=22)

    # 23. block 4 divider
    slide_section(prs, "Bayes in Action", page=23)

    # 24. sick friend
    slide_content(prs, "Sick friend",
                  [
                      ("Your 50-year-old friend has been a chain smoker since 18.",
                       {"size": 16, "color": DIM}),
                      ("There's a nasty cold going around.",
                       {"size": 16, "color": DIM}),
                      ("You go over to their house. You hear them cough.",
                       {"size": 16, "color": DIM, "space_after": 14}),
                      ("Three possibilities:",
                       {"size": 16, "bold": True, "color": WHITE,
                        "space_after": 8}),
                      ("Cold      ·      Stomach virus      ·      Lung cancer",
                       {"size": 18, "color": BLUE, "bold": True,
                        "space_after": 14}),
                      ("How likely is each?",
                       {"size": 16, "italic": True, "color": WHITE}),
                  ], page=24)

    # 25. first pass qualitative
    slide_content(prs, "First pass  —  qualitatively",
                  [
                      ("                                Prior            P(cough | h)        Prior × Likelihood",
                       {"size": 13, "bold": True, "color": MUTED,
                        "font": "Consolas", "space_after": 8}),
                      ("Cold                          high              high                       high",
                       {"size": 15, "color": GREEN, "bold": True,
                        "font": "Consolas"}),
                      ("Stomach virus           medium         low                          low",
                       {"size": 15, "color": DIM, "font": "Consolas"}),
                      ("Lung cancer              medium         high                         medium",
                       {"size": 15, "color": DIM, "font": "Consolas",
                        "space_after": 16}),
                      ("Before any math:  cold is most likely.",
                       {"size": 16, "bold": True, "color": WHITE}),
                  ], page=25)

    # 26. with numbers
    slide_content(prs, "With numbers",
                  [
                      ("                                P(h)         P(cough | h)        Numerator        Posterior",
                       {"size": 13, "bold": True, "color": MUTED,
                        "font": "Consolas", "space_after": 8}),
                      ("Cold                          0.45         0.9                       0.405                0.750",
                       {"size": 15, "color": GREEN, "bold": True,
                        "font": "Consolas"}),
                      ("Stomach virus           0.10         0.45                     0.045                0.083",
                       {"size": 15, "color": DIM, "font": "Consolas"}),
                      ("Lung cancer              0.10         0.9                       0.090                0.167",
                       {"size": 15, "color": DIM, "font": "Consolas",
                        "space_after": 14}),
                      ("Sum of numerators:  0.405 + 0.045 + 0.090  =  0.540.",
                       {"size": 14, "color": DIM}),
                      ("Divide each numerator by the sum to normalize.",
                       {"size": 14, "color": MUTED, "italic": True}),
                  ], page=26)

    # 27. pedagogical punchline
    slide_content(prs, "The pedagogical punchline",
                  [
                      ("P(cold | cough)  =  (0.9 · 0.45) / (0.9·0.45 + 0.45·0.10 + 0.9·0.10)",
                       {"size": 16, "color": WHITE, "font": "Cambria Math"}),
                      ("                  =  0.405 / 0.540  =  0.75",
                       {"size": 18, "color": BLUE, "bold": True,
                        "font": "Cambria Math", "space_after": 14}),
                      ("Even with a chain-smoking friend, P(lung cancer | cough) ≈ 17%  —  not dominant.",
                       {"size": 15, "color": WHITE, "space_after": 14}),
                      ("The cold wins because:",
                       {"size": 15, "bold": True, "color": WHITE,
                        "space_after": 6}),
                      ("1.  Priors matter.  Cold has a ~4.5× higher prior than cancer.",
                       {"size": 14, "color": DIM}),
                      ("2.  Likelihoods matter.  Cough is more likely from a cold (0.9) than a stomach virus (0.45).",
                       {"size": 14, "color": DIM}),
                      ("3.  Neither alone is enough.  You need the whole rule.",
                       {"size": 14, "color": DIM, "bold": True}),
                  ], page=27)

    # 28. K-T divider (or just header) — use small section divider style
    slide_section(prs, "The Kahneman–Tversky cab problem", page=28)

    # 29. setup
    slide_content(prs, "Setup",
                  [
                      ("•  A city has two cab companies:  Blue and Green.",
                       {"size": 16, "color": DIM}),
                      ("•  85% of cabs are Green.  15% are Blue.",
                       {"size": 16, "color": DIM}),
                      ("•  At night, a hit-and-run happens. A witness says the cab was Blue.",
                       {"size": 16, "color": DIM}),
                      ("•  The witness correctly identifies cab colors at night 80% of the time.",
                       {"size": 16, "color": DIM, "space_after": 18}),
                      ("What is the probability the cab was actually Blue?",
                       {"size": 18, "bold": True, "color": WHITE,
                        "space_after": 14}),
                      ("(Commit to a number before we compute.)",
                       {"size": 13, "color": MUTED, "italic": True}),
                  ], page=29)

    # 30. pass 2 area diagram
    slide_content(prs, "Pass 2  —  the area diagram",
                  [
                      ("Imagine 100 cabs.",
                       {"size": 16, "color": DIM, "space_after": 8}),
                      ("•  Blue cabs:  15.  Witness correctly says \"blue\" on 80% of them  →  12.",
                       {"size": 15, "color": DIM}),
                      ("•  Green cabs:  85.  Witness incorrectly says \"blue\" on 20% of them  →  17.",
                       {"size": 15, "color": DIM, "space_after": 10}),
                      ("Total \"blue\" reports:  12 + 17  =  29.",
                       {"size": 15, "color": DIM, "space_after": 12}),
                      ("P(Blue | witness says Blue)  =  12/29  ≈  0.41",
                       {"size": 22, "bold": True, "color": BLUE,
                        "font": "Cambria Math", "space_after": 16}),
                      ("Below 50%.  The witness is more likely wrong than right, despite being 80% accurate.",
                       {"size": 14, "bold": True, "color": WHITE}),
                  ], page=30)

    # 31. pass 3 formal
    slide_content(prs, "Pass 3  —  formal Bayes",
                  [
                      ("Hypotheses  h ∈ {B, G}.    Observation  o = \"witness says Blue.\"",
                       {"size": 15, "color": DIM, "space_after": 12}),
                      ("P(B | o)  =  P(o | B) · P(B)  /  [ P(o | B) · P(B)  +  P(o | G) · P(G) ]",
                       {"size": 16, "color": WHITE, "font": "Cambria Math",
                        "space_after": 8}),
                      ("           =  (0.80 · 0.15)  /  (0.80 · 0.15  +  0.20 · 0.85)",
                       {"size": 16, "color": WHITE, "font": "Cambria Math"}),
                      ("           =  0.12 / (0.12 + 0.17)  =  0.12 / 0.29  ≈  0.41",
                       {"size": 18, "bold": True, "color": BLUE,
                        "font": "Cambria Math", "space_after": 14}),
                      ("Same number as the area diagram.  The equation and the counting are the same calculation.",
                       {"size": 14, "color": MUTED, "italic": True}),
                  ], page=31)

    # 32. base-rate neglect
    slide_image_right(prs, "Base-rate neglect",
                      [
                          ("Most people report something close to 80%.",
                           {"size": 15, "color": DIM}),
                          ("They anchor on the likelihood and ignore the prior.",
                           {"size": 15, "color": DIM, "space_after": 14}),
                          ("Kahneman & Tversky (1972); Bar-Hillel (1980).",
                           {"size": 13, "color": MUTED, "italic": True,
                            "space_after": 12}),
                          ("Medical doctors fail this too, even for diagnoses",
                           {"size": 13, "color": MUTED, "italic": True}),
                          ("(Casscells, Schoenberger, Grayboys, 1978).",
                           {"size": 13, "color": MUTED, "italic": True}),
                      ],
                      IMG / "base_rate_plot_dark.png", page=32,
                      img_x=5.0, img_y=1.05, img_w=4.7, img_h=3.6)

    # 33. heuristics & biases
    slide_content(prs, "Heuristics and biases",
                  [
                      ("Heuristic.   A rule-of-thumb.  Usually works.  Cheap to apply.",
                       {"size": 16, "color": DIM, "space_after": 10}),
                      ("Bias.   Systematic error introduced when a heuristic is misapplied.",
                       {"size": 16, "color": DIM, "space_after": 16}),
                      ("Kahneman & Tversky argued that the mind is full of heuristics — good enough for most everyday inference but wrong in predictable ways when base rates are extreme, evidence is unfamiliar, or causal structure is unusual.",
                       {"size": 14, "color": DIM, "space_after": 12}),
                      ("This is a running theme for the semester.",
                       {"size": 14, "bold": True, "color": WHITE,
                        "italic": True}),
                  ], page=33)

    # 34. synthesis
    slide_content(prs, "Synthesis  —  the 5-step recipe",
                  [
                      ("This is the course's method.",
                       {"size": 16, "bold": True, "color": WHITE,
                        "space_after": 14}),
                      ("1.  Formalize the problem people face.",
                       {"size": 15, "color": DIM}),
                      ("2.  Formalize the knowledge they bring to it.",
                       {"size": 15, "color": DIM}),
                      ("3.  Apply Bayes' rule  —  compute what an ideal learner would infer.",
                       {"size": 15, "color": DIM}),
                      ("4.  Characterize the ideal learner's behavior.",
                       {"size": 15, "color": DIM}),
                      ("5.  Identify what knowledge and constraints must be assumed for model and human behavior to match.",
                       {"size": 15, "color": DIM, "space_after": 14}),
                      ("Every week we pick a different cognitive domain and run this loop.",
                       {"size": 14, "color": MUTED, "italic": True}),
                  ], page=34)

    # 35. block 5 divider
    slide_section(prs, "GenJAX, Admin, Homework", page=35)

    # 36. what is genjax
    slide_content(prs, "What is GenJAX?",
                  [
                      ("Bayes' rule is easy to write.  Enumerating hypothesis spaces by hand is not.",
                       {"size": 16, "color": DIM, "space_after": 10}),
                      ("For continuous distributions you can't enumerate at all.",
                       {"size": 16, "color": DIM, "space_after": 18}),
                      ("GenJAX is a probabilistic programming language.",
                       {"size": 17, "bold": True, "color": WHITE}),
                      ("You write the generative process as code. The machine does the enumeration and counting.",
                       {"size": 15, "color": DIM}),
                  ], page=36)

    # 37. genjax by example
    slide_content(prs, "GenJAX by example",
                  [
                      ("import genjax",
                       {"size": 14, "color": DIM, "font": "Consolas"}),
                      ("from genjax import gen, flip",
                       {"size": 14, "color": DIM, "font": "Consolas",
                        "space_after": 6}),
                      ("@gen",
                       {"size": 14, "color": GREEN, "font": "Consolas"}),
                      ("def chibany_day():",
                       {"size": 14, "color": DIM, "font": "Consolas"}),
                      ("    lunch  = flip(0.5) @ \"lunch\"",
                       {"size": 14, "color": DIM, "font": "Consolas"}),
                      ("    dinner = flip(0.5) @ \"dinner\"",
                       {"size": 14, "color": DIM, "font": "Consolas"}),
                      ("    return (lunch, dinner)",
                       {"size": 14, "color": DIM, "font": "Consolas",
                        "space_after": 12}),
                      ("•  @gen  —  this is a generative model, not a regular function.",
                       {"size": 13, "color": DIM}),
                      ("•  flip(0.5)  —  a Bernoulli random variable.",
                       {"size": 13, "color": DIM}),
                      ("•  @ \"lunch\"  —  name the random choice so we can condition on it later.",
                       {"size": 13, "color": DIM, "space_after": 8}),
                      ("This function IS the outcome space Ω from Block 3. Every call samples a point.",
                       {"size": 13, "bold": True, "color": WHITE,
                        "italic": True}),
                  ], page=37)

    # 38. preemptive callouts
    slide_content(prs, "Preemptive callouts",
                  [
                      ("•  Use  flip(p),  not  bernoulli(p).  bernoulli takes a logit, not a probability.",
                       {"size": 15, "color": DIM}),
                      ("•  The  @  operator here is not matrix multiplication. GenJAX overloads it.",
                       {"size": 15, "color": DIM}),
                      ("•  Output displays as  Array(0, dtype=int32).  Treat these as 0s and 1s.",
                       {"size": 15, "color": DIM}),
                      ("•  Everything runs in Google Colab.  No local installation.",
                       {"size": 15, "color": DIM}),
                  ], page=38)

    # 39. admin grading
    slide_grading(prs, page=39)

    # 40. admin rest
    slide_content(prs, "Admin  —  the rest",
                  [
                      ("•  Textbook:  A Narrative Introduction to Probability",
                       {"size": 15, "color": DIM}),
                      ("       https://josephausterweil.github.io/probintro/",
                       {"size": 13, "color": BLUE}),
                      ("•  Tooling:  Google Colab, GenJAX. No local setup.",
                       {"size": 15, "color": DIM}),
                      ("•  Office hours:  TBD — I'll poll.",
                       {"size": 15, "color": DIM}),
                      ("•  Email:  best-effort 36 h response (48 h on weekends).",
                       {"size": 15, "color": DIM}),
                      ("•  AI tools:  welcome as a technical resource; must cite; not for quiz/assignment answers.",
                       {"size": 15, "color": DIM}),
                      ("•  Full syllabus on the course website (coming up).",
                       {"size": 15, "color": DIM}),
                  ], page=40)

    # 41. readings overview 1
    slide_content(prs, "Semester readings  —  Weeks 3–7",
                  [
                      ("Each week includes 1-2 papers alongside the textbook. These may change.",
                       {"size": 13, "color": MUTED, "italic": True,
                        "space_after": 14}),
                      ("Wk 3  Conjugate Bayes & Topic Models",
                       {"size": 14, "bold": True, "color": WHITE}),
                      ("  Griffiths & Tenenbaum 2001 (Randomness)  ·  Steyvers & Griffiths 2000 (Topic Models)",
                       {"size": 12, "color": DIM, "space_after": 8}),
                      ("Wk 4  Generalization & Hierarchical Bayes",
                       {"size": 14, "bold": True, "color": WHITE}),
                      ("  Tenenbaum 1999 (Concept Learning)  ·  Xu & Tenenbaum 2000",
                       {"size": 12, "color": DIM, "space_after": 8}),
                      ("Wk 5  Hierarchical Bayes & Bayes Nets",
                       {"size": 14, "bold": True, "color": WHITE}),
                      ("  DeZoete 2019",
                       {"size": 12, "color": DIM, "space_after": 8}),
                      ("Wk 6  Causal Bayes Nets",
                       {"size": 14, "bold": True, "color": WHITE}),
                      ("  Lagnado 2002  ·  Griffiths et al. 2004 (Hidden Causes)  ·  Gerstenberg 2015",
                       {"size": 12, "color": DIM, "space_after": 8}),
                      ("Wk 7  Markov Chains & Networks",
                       {"size": 14, "bold": True, "color": WHITE}),
                      ("  Abbott et al. 2012 (Random Walks)  ·  Mukherjee 2017",
                       {"size": 12, "color": DIM}),
                  ], page=41)

    # 42. readings overview 2
    slide_content(prs, "Semester readings  —  Weeks 8–13",
                  [
                      ("Wk 8  Monte Carlo Methods",
                       {"size": 14, "bold": True, "color": WHITE}),
                      ("  Vul et al. 2009 (One and Done)  ·  Griffiths & Sanborn 2007",
                       {"size": 12, "color": DIM, "space_after": 8}),
                      ("Wk 9  SDT, MDPs & Reinforcement Learning",
                       {"size": 14, "bold": True, "color": WHITE}),
                      ("  Schultz 1997 (Dopamine & Reward)  ·  Daw et al. 2005  ·  Ho et al. 2015",
                       {"size": 12, "color": DIM, "space_after": 8}),
                      ("Wk 10  Inverse Reinforcement Learning",
                       {"size": 14, "bold": True, "color": WHITE}),
                      ("  Baker et al. 2007 (Theory of Mind)  ·  Ho et al. 2016, 2017",
                       {"size": 12, "color": DIM, "space_after": 8}),
                      ("Wk 11  Bayesian Nonparametrics",
                       {"size": 14, "bold": True, "color": WHITE}),
                      ("  Austerweil et al. 2015 (BNP Book Chapter)",
                       {"size": 12, "color": DIM, "space_after": 8}),
                      ("Wk 12  Deep Neural Networks",
                       {"size": 14, "bold": True, "color": WHITE}),
                      ("  LeCun et al. 2015 (Deep Learning)  ·  Lake et al. 2015  ·  Pereira et al. 2018",
                       {"size": 12, "color": DIM, "space_after": 8}),
                      ("Wk 13  Ethics & Adversarial ML",
                       {"size": 14, "bold": True, "color": WHITE}),
                      ("  Caliskan et al. 2017 (Bias)  ·  Buolamwini 2018  ·  Zhou & Firestone 2019",
                       {"size": 12, "color": DIM}),
                  ], page=42)

    # 43. homework
    slide_content(prs, "Homework for Week 2",
                  [
                      ("1.  Textbook T1 Ch 1-3  —  reinforces everything from Block 3.",
                       {"size": 15, "color": DIM, "bold": True}),
                      ("       intro/01_goals.md  ·  02_hungry.md  ·  03_prob_count.md",
                       {"size": 13, "color": MUTED, "font": "Consolas",
                        "space_after": 8}),
                      ("2.  Textbook T2 Ch 0-1  —  gets GenJAX running in Colab.",
                       {"size": 15, "color": DIM, "bold": True}),
                      ("       genjax/00_getting_started.md  ·  01_python_basics.md",
                       {"size": 13, "color": MUTED, "font": "Consolas",
                        "space_after": 8}),
                      ("3.  Start thinking about a final-project topic. We'll talk Week 2.",
                       {"size": 15, "color": DIM, "bold": True,
                        "space_after": 12}),
                      ("Week 2 has a short optional self-check quiz (Intro Probability Theory 1).",
                       {"size": 13, "color": MUTED, "italic": True}),
                  ], page=43)

    # 44. questions
    slide_centered(prs, "Questions?", "Thanks — see you next Friday.",
                   page=44)

    out = HERE / "week01.pptx"
    prs.save(str(out))
    print(f"Wrote {out} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
